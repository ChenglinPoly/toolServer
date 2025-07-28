import sys
import asyncio
import subprocess
import shutil
import git
from pathlib import Path
from typing import List, Optional
from .base_tool import LocalTool
from utils.response import ToolResponse
from utils.logger import global_logger


class ExecuteCodeTool(LocalTool):
    """Python代码执行工具"""
    
    def __init__(self):
        super().__init__()
        self.tool_name = "execute_code"
        self.description = "执行Python代码文件"
    
    async def execute(self, task_id: str, workspace_path: Path, file_path: str, timeout: int = 300, **kwargs) -> ToolResponse:
        try:
            if not file_path:
                return ToolResponse(success=False, error="file_path is required")
            
            task_path = self.get_task_path(task_id, workspace_path)
            full_path = task_path / file_path
            
            if not full_path.exists():
                return ToolResponse(success=False, error=f"File not found: {file_path}")
            
            # 使用任务的虚拟环境
            venv_python = task_path / "code_env" / "bin" / "python"
            if not venv_python.exists():
                venv_python = sys.executable
            
            # 切换到代码运行目录
            code_run_dir = task_path / "code_run"
            code_run_dir.mkdir(exist_ok=True)
            
            try:
                process = await asyncio.create_subprocess_exec(
                    str(venv_python), str(full_path),
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    cwd=str(code_run_dir)
                )
                
                stdout, stderr = await asyncio.wait_for(
                    process.communicate(), 
                    timeout=timeout
                )
                
                return ToolResponse(
                    success=True,
                    data={
                        "exit_code": process.returncode,
                        "stdout": stdout.decode('utf-8') if stdout else "",
                        "stderr": stderr.decode('utf-8') if stderr else "",
                        "executed_file": file_path
                    }
                )
                
            except asyncio.TimeoutError:
                try:
                    process.kill()
                    await process.wait()
                except:
                    pass
                return ToolResponse(
                    success=False,
                    error=f"Execution timeout after {timeout} seconds"
                )
                
        except Exception as e:
            return ToolResponse(success=False, error=str(e))


class ExecuteShellTool(LocalTool):
    """Shell命令执行工具"""
    
    def __init__(self):
        super().__init__()
        self.tool_name = "execute_shell"
        self.description = "执行Shell命令"
    
    async def execute(self, task_id: str, workspace_path: Path, command: str, timeout: int = 60, workdir: str = '', **kwargs) -> ToolResponse:
        try:
            if not command:
                return ToolResponse(success=False, error="command is required")
            
            task_path = self.get_task_path(task_id, workspace_path)
            if workdir:
                work_path = task_path / workdir
            else:
                work_path = task_path / "code_run"
            
            work_path.mkdir(parents=True, exist_ok=True)
            
            try:
                process = await asyncio.create_subprocess_shell(
                    command,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    cwd=str(work_path)
                )
                
                stdout, stderr = await asyncio.wait_for(
                    process.communicate(), 
                    timeout=timeout
                )
                
                return ToolResponse(
                    success=True,
                    data={
                        "exit_code": process.returncode,
                        "stdout": stdout.decode('utf-8') if stdout else "",
                        "stderr": stderr.decode('utf-8') if stderr else "",
                        "command": command
                    }
                )
                
            except asyncio.TimeoutError:
                try:
                    process.kill()
                    await process.wait()
                except:
                    pass
                return ToolResponse(
                    success=False,
                    error=f"Command timeout after {timeout} seconds"
                )
                
        except Exception as e:
            return ToolResponse(success=False, error=str(e))


class PipInstallTool(LocalTool):
    """Python包安装工具"""
    
    def __init__(self):
        super().__init__()
        self.tool_name = "pip_install"
        self.description = "在虚拟环境中安装Python包"
    
    async def execute(self, task_id: str, workspace_path: Path, packages: List[str] = None, **kwargs) -> ToolResponse:
        try:
            if not packages:
                return ToolResponse(success=False, error="packages is required")
            
            task_path = self.get_task_path(task_id, workspace_path)
            pip_path = task_path / "code_env" / "bin" / "pip"
            
            # 如果虚拟环境不存在，创建它
            if not pip_path.exists():
                venv_path = task_path / "code_env"
                if venv_path.exists():
                    shutil.rmtree(venv_path)
                
                process = await asyncio.create_subprocess_exec(
                    sys.executable, "-m", "venv", str(venv_path),
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                stdout, stderr = await process.communicate()
                
                if process.returncode != 0:
                    return ToolResponse(
                        success=False,
                        error=f"Failed to create virtual environment: {stderr.decode('utf-8')}"
                    )
                
                if not pip_path.exists():
                    return ToolResponse(
                        success=False,
                        error=f"pip not found after creating virtual environment: {pip_path}"
                    )
            
            # 安装包
            if isinstance(packages, str):
                packages = [packages]
            
            results = []
            for package in packages:
                process = await asyncio.create_subprocess_exec(
                    str(pip_path), "install", package,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                stdout, stderr = await process.communicate()
                
                results.append({
                    "package": package,
                    "success": process.returncode == 0,
                    "stdout": stdout.decode('utf-8') if stdout else "",
                    "stderr": stderr.decode('utf-8') if stderr else ""
                })
            
            all_success = all(r["success"] for r in results)
            
            return ToolResponse(
                success=all_success,
                data={"results": results}
            )
            
        except Exception as e:
            return ToolResponse(success=False, error=str(e))


class GitCloneTool(LocalTool):
    """Git仓库克隆工具"""
    
    def __init__(self):
        super().__init__()
        self.tool_name = "git_clone"
        self.description = "克隆Git仓库"
    
    async def execute(self, task_id: str, workspace_path: Path, repo_url: str, target_dir: str = '', branch: str = None, token: str = None, **kwargs) -> ToolResponse:
        try:
            if not repo_url:
                return ToolResponse(success=False, error="repo_url is required")
            
            task_path = self.get_task_path(task_id, workspace_path)
            upload_dir = task_path / "upload"
            
            if target_dir:
                clone_path = upload_dir / target_dir
            else:
                repo_name = repo_url.rstrip('/').split('/')[-1]
                if repo_name.endswith('.git'):
                    repo_name = repo_name[:-4]
                clone_path = upload_dir / repo_name
            
            clone_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 处理认证
            original_repo_url = repo_url
            if token and 'github.com' in repo_url:
                repo_url = repo_url.replace('https://', f'https://{token}@')
            
            # 构建git命令
            git_cmd = ["git", "clone"]
            if branch:
                git_cmd.extend(["-b", branch])
            git_cmd.extend([repo_url, str(clone_path)])
            
            try:
                process = await asyncio.create_subprocess_exec(
                    *git_cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                stdout, stderr = await process.communicate()
                
                if process.returncode != 0:
                    return ToolResponse(
                        success=False,
                        error=f"Git clone failed: {stderr.decode('utf-8')}"
                    )
                
                # 获取仓库信息
                async def get_git_info(cmd_args):
                    proc = await asyncio.create_subprocess_exec(
                        "git", *cmd_args,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                        cwd=str(clone_path)
                    )
                    stdout, _ = await proc.communicate()
                    return stdout.decode('utf-8').strip() if proc.returncode == 0 else ""
                
                branch_name, commit_hash, commit_message, commit_author, commit_date = await asyncio.gather(
                    get_git_info(["rev-parse", "--abbrev-ref", "HEAD"]),
                    get_git_info(["rev-parse", "HEAD"]),
                    get_git_info(["log", "-1", "--pretty=format:%s"]),
                    get_git_info(["log", "-1", "--pretty=format:%an <%ae>"]),
                    get_git_info(["log", "-1", "--pretty=format:%ci"]),
                    return_exceptions=True
                )
                
                info = {
                    "clone_path": str(clone_path.relative_to(task_path)),
                    "repo_url": original_repo_url,
                    "branch": branch_name if isinstance(branch_name, str) else "unknown",
                    "commit": commit_hash if isinstance(commit_hash, str) else "unknown",
                    "commit_message": commit_message if isinstance(commit_message, str) else "unknown",
                    "commit_author": commit_author if isinstance(commit_author, str) else "unknown",
                    "commit_date": commit_date if isinstance(commit_date, str) else "unknown"
                }
                
                return ToolResponse(success=True, data=info)
                
            except asyncio.TimeoutError:
                return ToolResponse(
                    success=False,
                    error="Git clone operation timed out"
                )
            
        except Exception as e:
            return ToolResponse(success=False, error=str(e))


class ParseDocumentTool(LocalTool):
    """文档解析工具"""
    
    def __init__(self):
        super().__init__()
        self.tool_name = "parse_document"
        self.description = "解析文档（PDF、Word、PPT、Markdown）"
    
    async def execute(self, task_id: str, workspace_path: Path, file_path: str, **kwargs) -> ToolResponse:
        try:
            if not file_path:
                return ToolResponse(success=False, error="file_path is required")
            
            task_path = self.get_task_path(task_id, workspace_path)
            full_path = task_path / file_path
            
            if not full_path.exists():
                return ToolResponse(success=False, error=f"File not found: {file_path}")
            
            ext = full_path.suffix.lower()
            content = ""
            metadata = {"file_type": ext}
            
            try:
                if ext == '.pdf':
                    import PyPDF2
                    with open(full_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        metadata["pages"] = len(reader.pages)
                        content = ""
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                            
                elif ext in ['.docx', '.doc']:
                    from docx import Document
                    doc = Document(full_path)
                    paragraphs = []
                    for para in doc.paragraphs:
                        if para.text.strip():
                            paragraphs.append(para.text)
                    content = "\n".join(paragraphs)
                    metadata["paragraphs"] = len(paragraphs)
                    
                elif ext in ['.pptx', '.ppt']:
                    from pptx import Presentation
                    prs = Presentation(full_path)
                    slides_text = []
                    for slide in prs.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                        if slide_text:
                            slides_text.append("\n".join(slide_text))
                    content = "\n\n---\n\n".join(slides_text)
                    metadata["slides"] = len(prs.slides)
                    
                elif ext in ['.md', '.markdown']:
                    with open(full_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    metadata["lines"] = len(content.splitlines())
                    
                else:
                    return ToolResponse(
                        success=False, 
                        error=f"Unsupported file type: {ext}"
                    )
                
                return ToolResponse(
                    success=True,
                    data={
                        "content": content,
                        "metadata": metadata,
                        "file_path": file_path,
                        "char_count": len(content),
                        "word_count": len(content.split())
                    }
                )
                
            except Exception as e:
                return ToolResponse(
                    success=False, 
                    error=f"Failed to parse {ext} file: {str(e)}"
                )
                
        except Exception as e:
            return ToolResponse(success=False, error=str(e)) 